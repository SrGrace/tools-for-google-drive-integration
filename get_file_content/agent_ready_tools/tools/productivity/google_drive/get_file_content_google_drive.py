
from pathlib import Path
import sys

test_dir = Path(__file__).parent
BASE_DIR = 'agent_ready_tools'
MAX_DEPTH = 10

while test_dir.name != BASE_DIR:
    test_dir = test_dir.parent
    MAX_DEPTH -= 1
    if MAX_DEPTH == 0:
        raise RecursionError(f"'{BASE_DIR}' not found in path: {__file__}")
parent_path = test_dir.parent.resolve()

sys.path.append(str(parent_path))
    
from typing import Optional

from ibm_watsonx_orchestrate.agent_builder.tools import tool
from pydantic.dataclasses import dataclass

from agent_ready_tools.clients.google_client import get_google_client
from agent_ready_tools.utils.tool_credentials import GOOGLE_CONNECTIONS
import base64
import mammoth
import io
import openpyxl
import PyPDF2
from pptx import Presentation

@dataclass
class FileContentGoogleDriveResult:
    """Retrieves the contents of a file in Google Drive."""
    content: str


@tool(expected_credentials=GOOGLE_CONNECTIONS)
def get_file_content_google_drive(
    file_id: str,
    alt: str = "media",
    file_type: Optional[str] = None,
) -> FileContentGoogleDriveResult:
    """
    Retrieves the content of a file in Google Drive.

    Args:
        file_id: The id of the file returned by the `get_files` tool.
        alt: Returns the actual file content.
        file_type: The type of the file returned by the `get_files` tool.

    Returns:
        The content of a file.
    """

    client = get_google_client()
    
    # Determine export MIME type for Google native formats
    export_mime_type = None
    if file_type == "application/vnd.google-apps.document":
        export_mime_type = "text/plain"
    elif file_type == "application/vnd.google-apps.spreadsheet":
        export_mime_type = "text/csv"
    elif file_type == "application/vnd.google-apps.presentation":
        export_mime_type = "text/plain"
    
    # Download the file
    download_result = client.download_file(file_id=file_id, export_mime_type=export_mime_type)
    content = download_result["content"]
    
    # If content is base64-encoded, try to extract text based on file type
    try:
        # Check if content is base64 (will fail decode if it's already text)
        file_bytes = base64.b64decode(content)
        
        # Handle .docx files
        if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
            content = result.value
        
        # Handle .pptx files (PowerPoint)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(file_bytes))
            slides_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slides_content.append(f"\n=== Slide {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides_content.append(shape.text)
            content = "\n".join(slides_content)
        
        # Handle .xlsx files (Excel)
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
            sheets_content = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheets_content.append(f"\n=== Sheet: {sheet_name} ===\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    sheets_content.append(row_text)
            content = "\n".join(sheets_content)
        
        # Handle PDFs
        elif file_type == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            pages_content = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                pages_content.append(f"\n=== Page {page_num} ===\n")
                pages_content.append(page.extract_text())
            content = "\n".join(pages_content)
        
        # For other binary types, keep as base64
        else:
            content = f"[Binary file of type {file_type}. Content is base64 encoded.]\n{content[:200]}..."
            
    except Exception as e:
        # If any processing fails or content is already text, keep as-is
        pass
    
    return FileContentGoogleDriveResult(content=content)

"""
orchestrate tools import -k python -p . -f agent_ready_tools/tools/productivity/google_drive/get_file_content_google_drive.py -a google_key_value_ibm_184bdbd3
"""